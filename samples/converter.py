from pyth.plugins.plaintext.writer import PlaintextWriter
from pyth.plugins.rtf15.reader import Rtf15Reader 
from pyth import document 

from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter 
from pdfminer.converter import TextConverter 
from pdfminer.layout import LAParams 
from pdfminer.pdfpage import PDFPage 

from reportlab.platypus import SimpleDocTemplate, Table, TableStyle 
from reportlab.pdfgen.canvas import Canvas
from reportlab.lib.pagesizes import letter 
from reportlab.lib.units import inch 
from reportlab.lib import colors 

from pptx import Presentation 
from pptx.util import Inches

from docx import Document 
from io import StringIO 
from io import BytesIO 
from PIL import Image 
import pandas as pd
import sys, csv, os  
import pypandoc

class Documents():
    
    class CSV():
        
        def __init__(self, filename):
            self.filename = filename 
        
        def ToDOCX(self, outputfile):
            doc = Document() 
            
            with open(self.filename, 'r') as csvfile:
                reader = csv.reader(csvfile)
                for row in reader:
                    doc.add_paragraph(', '.join(row))
            
            doc.save(outputfile)
            return True
        
        def ToHTML(self, outputfile):
            df = pd.read_csv(self.filename)
            html = df.to_html()
            
            with open(outputfile, 'w') as file:
                file.write(html)
            return True 
        
        def ToPDF(self, outputfile):
            df = pd.read_csv(self.filename)
            c = Canvas(outputfile, pagesize = letter)
            data = [df.columns.tolist()] + df.values.tolist()
            table = Table(data)
            
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'), 
                ('FONTSIZE', (0, 0), (-1, 0), 14),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                ('ALIGN', (0, 1), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 12),
                ('BOTTOMPADDING', (0, 1), (-1, -1), 10),
            ])
            
            table.setStyle(style)
            elements = [] 
            elements.append(table) 
            
            doc = SimpleDocTemplate(outputfile, pagesize = letter)
            doc.build(elements)
            c.save()
            
            return True 
        
        def ToRTF(self, outputfile):
            df = pd.read_csv(self.filename)
            table = document.Table([document.Row([document.Text('') for j in range(len(df.columns))]) for i in range(len(df.index) + 1)])

            for i, col in enumerate(df.columns):
                table[0][i].content = col 
            
            for i, row in enumerate(df.values):
                for j, cell, in enumerate(row):
                    table[i + 1][j].content = str(cell)
            
            doc = document.Document()
            doc.append(table)
            
            with open(outputfile, 'wb') as file:
                file.write(Rtf15Reader.write(doc).getvalue())
            return True 

        def ToTXT(self, outputfile):
            csvfile = open(self.filename, 'r')
            txtfile = open(outputfile, 'w')
            reader = csv.reader(csvfile)
            
            for row in reader:
                txtfile.write('\t'.join(row) + '\n')
            
            csvfile.close()
            txtfile.close()
            
            return True 
        
        def CSVToXLSX(self, outputfile):
            df = pd.read_csv(self.filename)
            df.to_excel(outputfile, index = False)
            return True 

        def ToPPTX(self, outputfile):
            with open(self.filename, 'r') as csvfile:
                reader = csv.reader(csvfile)
                rows = list(reader)
            prs = Presentation() 
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            num_rows = len(rows)
            num_cols = len(rows[0])
            table = slide.shapes.add_table(num_rows, num_cols, Inches(1.0), Inches(2.0), Inches(num_cols * 2), Inches(num_rows * 2), Inches(num_rows * 0.5)).table

            for i, row in enumerate(rows):
                for j, value in enumerate(row):
                    table.cell(i, j).text = str(value)
            prs.save(outputfile)
            
            return True 

class Images():
    
    class JIF():
        
        def __init__(self, filename):
            self.filename = filename 
            
        def ToJFIF(self, outputfile):
            with Image.open(self.filename) as img:
                img.save(outputfile, 'JPEG', quality = 95, jfif_version = (1, 1))
            return True 
        
        def ToJPG(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGB')
                img.save(outputfile, 'JPEG', quality = 95)
            return True 
        
        def ToPNG(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGBA')
                img.save(outputfile, 'PNG')
            return True 
        
        def ToPSD(self, outputfile):
            with Image.open(self.filename) as img:
                img.save(outputfile, 'PSD')
            return True 
        
        def ToTIFF(self, outputfile):
            with Image.open(self.filename) as img:
                img.save(outputfile, 'TIFF')
            return True 
        
        def ToWEBP(self, outputfile):
            with Image.open(self.filename) as img:
                img.save(outputfile, 'WEBP', quality = 95)
            return True 
    
    class JFIF():
        
        def __init__(self, filename):
            self.filename = filename 
            
        def ToJIF(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGB')
                img.save(outputfile, 'JIF')
            return True 
        
        def ToJPG(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGB')
                img.save(outputfile, 'JPEG')
            return True 
        
        def ToPNG(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGBA')
                img.save(outputfile, 'PNG')
            return True 
        
        def ToPSD(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGBA')
                img.save(outputfile, 'PSD')
            return True 

        def ToSVG(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGBA')
                img.save(outputfile, 'SVG')
            return True 
        
        def ToTIFF(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGBA')
                img.save(outputfile, 'TIFF')
            return True 
        
        def ToWEBP(self, outputfile):
            with Image.open(self.filename) as img:
                img = img.convert('RGBA')
                img.save(outputfile, 'WEBP')
            return True 

# PPTX & PDF 
def PPTXToPDF(inputfile, outputfile):
    prs = Presentation(inputfile)
    prs.save(outputfile, 'pdf')
    return True 

def PDFToPPTX(inputfile, outputfile):
    manager = PDFResourceManager() 
    textio = StringIO()
    laparams = LAParams() 
    converter = TextConverter(manager, textio, laparams = laparams)
    interpreter = PDFPageInterpreter(manager, converter)
    
    with open(inputfile, 'rb') as file:
        for page in PDFPage.get_pages(file):
            interpreter.process_page(page)
    text = textio.getvalue() 
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    text_box = slide.shapes.add_textbox(left = 0, top = 0, width = prs.slide_width, height = prs.slide_height)
    text_frame = text_box.text_frame 
    text_frame.clear() 
    text_frame.text = text 
    
    prs.save(outputfile)
    return True 

# XLSX & CSV 
def XLSXToCSV(inputfile, outputfile):
    df = pd.read_excel(inputfile)
    df.to_csv(outputfile, index = False)
    return True 

# DOCX & PDF 
def DOCXToPDF(inputfile, outputfile):
    doc = Document(inputfile)
    pdffile = BytesIO()
    canvas_obj = Canvas(pdffile)
    
    for para in doc.paragraphs:
        canvas_obj.drawString(0, 0, para.text)
    canvas_obj.save()
    
    with open(outputfile, 'wb') as file:
        file.write(pdffile.getvalue())
        file.close()
    return True 

def PDFToDOCX(inputfile, outputfile):
    pypandoc.convert_file(inputfile,
    'docx', outputfile = outputfile)
    return True